#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Video2PPT - Convert videos to PowerPoint presentations
Extract key frames from videos and generate PPT
"""

import os
import sys
import cv2
import argparse
import logging
from pathlib import Path
from typing import List
import numpy as np

try:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches
    import imagehash
except ImportError as e:
    print(f"Error: Missing required library - {e}")
    print("Please run: pip install -r requirements.txt")
    sys.exit(1)


# Configure logging
LOG_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "video2ppt.log")

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(LOG_FILE, encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)


class Video2PPT:
    """Video to PPT converter"""

    def __init__(self, video_path: str, output_path: str = None, fps_interval: int = 1):
        """
        Initialize converter

        Args:
            video_path: Path to input video file
            output_path: Path to output PPT file
            fps_interval: Extract one frame every N seconds
        """
        self.video_path = video_path
        self.fps_interval = fps_interval
        self.frames: List[str] = []  # Store frame paths
        self.frames_dir = None
        self.hash_threshold = 15  # 哈希差异阈值，>15保留，≤15丢弃
        self.pixel_threshold = 30  # 像素差异阈值，>30直接保存，≤30走pHash
        self.last_hash = None     # 上一保留帧的哈希

        if not os.path.exists(video_path):
            raise FileNotFoundError(f"Video file not found: {video_path}")

        if output_path is None:
            base_name = Path(video_path).stem
            output_path = f"{base_name}_output.pptx"

        self.output_path = output_path
        logger.info(f"Initializing converter: {video_path} -> {output_path}")

    def _compute_hash_from_array(self, frame) -> imagehash.ImageHash:
        """从numpy数组计算感知哈希（不写盘）"""
        try:
            # BGR→RGB + PIL Image
            rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            img = Image.fromarray(rgb)
            return imagehash.phash(img)
        except Exception as e:
            logger.warning(f"Failed to compute hash from array: {e}")
            return None

    def extract_frames(self) -> List[str]:
        """提帧阶段：先算pHash，通过再写盘，避免无用I/O"""
        logger.info("Starting frame extraction...")

        video_stem = Path(self.video_path).stem
        self.frames_dir = f"{video_stem}_frames"
        os.makedirs(self.frames_dir, exist_ok=True)

        cap = cv2.VideoCapture(self.video_path)

        if not cap.isOpened():
            raise ValueError(f"Unable to open video file: {self.video_path}")

        fps = cap.get(cv2.CAP_PROP_FPS)
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = total_frames / fps if fps > 0 else 0

        logger.info(f"Video info - FPS: {fps:.2f}, Total frames: {total_frames}, Duration: {duration:.2f}s")

        frame_interval = int(fps * self.fps_interval)
        frame_count = 0
        extracted_count = 0
        last_frame = None  # 上一帧的numpy数组

        while True:
            ret, frame = cap.read()

            if not ret:
                break

            if frame_count % frame_interval == 0:
                frame_path = os.path.join(self.frames_dir, f"frame_{extracted_count:04d}.jpg")

                if last_frame is not None:
                    # 快速像素差异比较
                    diff = np.mean(np.abs(frame.astype(float) - last_frame.astype(float)))

                    if diff > self.pixel_threshold:
                        # 差异大，直接算pHash，通过再写盘
                        current_hash = self._compute_hash_from_array(frame)
                        if current_hash is None or self.last_hash is None or \
                           abs(current_hash - self.last_hash) > self.hash_threshold:
                            cv2.imwrite(frame_path, frame)
                            self.frames.append(frame_path)
                            self.last_hash = current_hash if current_hash is not None else None
                            extracted_count += 1
                            logger.debug(f"Frame {extracted_count}: diff={diff:.2f}, direct save")
                        else:
                            logger.debug(f"Frame skipped (similar): diff={diff:.2f}")
                    else:
                        # 疑似相似，先在内存中算pHash精确比较
                        current_hash = self._compute_hash_from_array(frame)
                        if current_hash is not None:
                            if self.last_hash is None or \
                               abs(current_hash - self.last_hash) > self.hash_threshold:
                                cv2.imwrite(frame_path, frame)
                                self.frames.append(frame_path)
                                self.last_hash = current_hash
                                extracted_count += 1
                                logger.debug(f"Frame {extracted_count}: diff={diff:.2f}, pHash save")
                            else:
                                # 相似，跳过
                                logger.debug(f"Frame skipped (similar): diff={diff:.2f}")
                        else:
                            # pHash计算失败，保存帧
                            cv2.imwrite(frame_path, frame)
                            self.frames.append(frame_path)
                            self.last_hash = None
                            extracted_count += 1
                else:
                    # 第一帧，直接保存
                    cv2.imwrite(frame_path, frame)
                    self.frames.append(frame_path)
                    self.last_hash = self._compute_hash_from_array(frame)
                    extracted_count += 1

                last_frame = frame

                if extracted_count % 10 == 0:
                    logger.info(f"Extracted {extracted_count} frames")

            frame_count += 1

        cap.release()
        logger.info(f"Frame extraction complete. Total frames extracted: {extracted_count}")
        return self.frames

    def generate_ppt(self) -> None:
        """Generate PowerPoint presentation"""
        logger.info("Starting PowerPoint generation...")

        if not self.frames:
            logger.error("No frame data available")
            return

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Add frame slides
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout

        for idx, frame_path in enumerate(self.frames, 1):
            logger.info(f"Processing slide {idx}/{len(self.frames)}...")

            slide = prs.slides.add_slide(blank_slide_layout)

            # Add frame image, fill entire slide (edge-aligned)
            left = Inches(0)
            top = Inches(0)
            width = Inches(10)     # Slide width
            height = Inches(7.5)   # Slide height
            pic = slide.shapes.add_picture(frame_path, left, top, width=width, height=height)

        # Save presentation
        prs.save(self.output_path)
        logger.info(f"PowerPoint saved: {self.output_path}")

    def cleanup(self) -> None:
        """Clean up temporary files"""
        if self.frames_dir and os.path.exists(self.frames_dir):
            import shutil
            shutil.rmtree(self.frames_dir)
            logger.info("Temporary files cleaned up")

    def convert(self) -> None:
        """完整转换流程：提帧(含去重) -> 生成PPT"""
        try:
            self.extract_frames()
            self.generate_ppt()
            logger.info("Conversion completed successfully!")
        except Exception as e:
            logger.error(f"Error during conversion: {e}")
            raise
        finally:
            self.cleanup()


def main():
    """Main function"""
    parser = argparse.ArgumentParser(
        description='Convert video files to PowerPoint presentations',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python main.py input_video.mp4
  python main.py input_video.mp4 -o output.pptx
  python main.py input_video.mp4 -i 2  (extract one frame every 2 seconds)
        """
    )

    parser.add_argument('video', help='Path to input video file')
    parser.add_argument('-o', '--output', help='Path to output PPT file (default: input_output.pptx)')
    parser.add_argument('-i', '--interval', type=int, default=1,
                       help='Frame extraction interval in seconds (default: 1)')

    args = parser.parse_args()

    try:
        converter = Video2PPT(args.video, args.output, args.interval)
        converter.convert()
    except Exception as e:
        logger.error(f"Error: {e}")
        sys.exit(1)


if __name__ == '__main__':
    main()
